#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
하루온 (HaruOn) 발표 PPT 생성기
Style : Simple IR Deck — mint/teal accent, white/light-gray background
Output: HaruOn_PPT.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ═══════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════
SW = Inches(13.33)
SH = Inches(7.5)

TEAL   = RGBColor(0x00, 0xC9, 0xA7)
TEAL_D = RGBColor(0x00, 0x9E, 0x85)
BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xF3, 0xF3, 0xF3)
GRAY   = RGBColor(0x77, 0x77, 0x77)
LGRAY  = RGBColor(0xD0, 0xD0, 0xD0)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
BADGE  = RGBColor(0x28, 0x28, 0x28)

FN = "맑은 고딕"

SHP_RECT  = 1
SHP_RRECT = 5
SHP_OVAL  = 9

# ═══════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════
def new_prs():
    p = Presentation()
    p.slide_width  = SW
    p.slide_height = SH
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(sl, col=BG):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = col

def _fl(s, fill, line, lw):
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = lw

def R(sl, l, t, w, h, fill=None, line=None, lw=Pt(0.75)):
    s = sl.shapes.add_shape(SHP_RECT, l, t, w, h)
    _fl(s, fill, line, lw)
    return s

def RR(sl, l, t, w, h, fill=None, line=None, lw=Pt(0.75)):
    s = sl.shapes.add_shape(SHP_RRECT, l, t, w, h)
    _fl(s, fill, line, lw)
    return s

def OV(sl, l, t, w, h, fill=None, line=None, lw=Pt(0.75)):
    s = sl.shapes.add_shape(SHP_OVAL, l, t, w, h)
    _fl(s, fill, line, lw)
    return s

def _round_max(s):
    try:
        for gd in s._element.iter(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}gd'
        ):
            gd.set('fmla', 'val 50000')
    except Exception:
        pass

def T(sl, l, t, w, h, text='', size=18, bold=False,
      color=BLACK, align=PP_ALIGN.LEFT, wrap=True,
      anchor=MSO_ANCHOR.TOP, italic=False):
    tx = sl.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FN
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tx

def ML(sl, l, t, w, h, items, default_size=16,
       default_color=BLACK, align=PP_ALIGN.LEFT,
       wrap=True, anchor=MSO_ANCHOR.TOP, spc=0):
    """Multi-line textbox.
    items: list of str  |  (text, size, bold, color)  |  dict
    """
    tx = sl.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, item in enumerate(items):
        if isinstance(item, str):
            item = (item, default_size, False, default_color)
        elif isinstance(item, dict):
            item = (
                item.get('text', ''),
                item.get('size', default_size),
                item.get('bold', False),
                item.get('color', default_color),
            )
        text, sz, bld, col = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spc and i > 0:
            p.space_before = Pt(spc)
        r = p.add_run()
        r.text  = text
        r.font.name  = FN
        r.font.size  = Pt(sz)
        r.font.bold  = bld
        r.font.color.rgb = col
    return tx

def circ(sl, cx, cy, r, fill=TEAL, text='', tc=WHITE, ts=18, tb=True):
    s = OV(sl, cx - r, cy - r, r * 2, r * 2, fill=fill)
    if text:
        tf = s.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(text)
        run.font.name  = FN
        run.font.size  = Pt(ts)
        run.font.bold  = tb
        run.font.color.rgb = tc
    return s

def pill(sl, l, t, w, h, fill=TEAL, text='', tc=WHITE, ts=12, tb=True):
    s = RR(sl, l, t, w, h, fill=fill)
    _round_max(s)
    tf = s.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name  = FN
    run.font.size  = Pt(ts)
    run.font.bold  = tb
    run.font.color.rgb = tc
    return s

def hdr(sl, en, ko='', offset_y=Inches(0)):
    """Standard section header (teal bar + EN title + KO subtitle)."""
    R(sl, 0, 0, SW, Inches(0.06), fill=TEAL)
    T(sl,
      l=0, t=Inches(0.2) + offset_y, w=SW, h=Inches(1.0),
      text=en, size=44, bold=True, color=BLACK,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if ko:
        T(sl,
          l=0, t=Inches(1.25) + offset_y, w=SW, h=Inches(0.4),
          text=ko, size=14, color=GRAY,
          align=PP_ALIGN.CENTER)

def card(sl, l, t, w, h, fill=WHITE, line=LGRAY, radius=True):
    if radius:
        s = RR(sl, l, t, w, h, fill=fill, line=line, lw=Pt(0.5))
    else:
        s = R(sl, l, t, w, h, fill=fill, line=line, lw=Pt(0.5))
    return s

# ═══════════════════════════════════════════
# SLIDE 01 — COVER
# ═══════════════════════════════════════════
def s01_cover(prs):
    sl = blank(prs)
    bg(sl, RGBColor(0xEB, 0xEB, 0xEB))

    # Decorative teal arcs (large offset ovals — stroke only)
    # Top-right arc
    OV(sl,
       l=Inches(8.5), t=Inches(-1.0),
       w=Inches(6.5), h=Inches(6.5),
       fill=None, line=TEAL, lw=Pt(8))
    # Bottom-left arc
    OV(sl,
       l=Inches(-1.5), t=Inches(2.5),
       w=Inches(6), h=Inches(6),
       fill=None, line=LGRAY, lw=Pt(8))

    # Center card (white rounded rectangle)
    card(sl,
         l=Inches(2.8), t=Inches(1.1),
         w=Inches(7.7), h=Inches(5.3),
         fill=WHITE, line=LGRAY)

    # "SIMPLE STYLE" small label
    T(sl, l=Inches(2.8), t=Inches(1.7), w=Inches(7.7), h=Inches(0.4),
      text='MINI PROJECT', size=13, bold=False, color=GRAY,
      align=PP_ALIGN.CENTER)

    # Main title
    ML(sl,
       l=Inches(2.8), t=Inches(2.15), w=Inches(7.7), h=Inches(1.5),
       items=[
           ('하루온', 52, True, TEAL),
           ('HaruOn Planner', 36, True, BLACK),
       ],
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Subtitle
    T(sl,
      l=Inches(2.8), t=Inches(3.7), w=Inches(7.7), h=Inches(0.6),
      text='연간 계획부터 10분 실행까지, AI와 함께 채우는 플래너',
      size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # Team info
    T(sl,
      l=Inches(2.8), t=Inches(4.7), w=Inches(7.7), h=Inches(0.4),
      text='KDT WebAI 4팀  |  김태동 · 박성훈 · 이현아 · 이창민  |  2026.05',
      size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 02 — TABLE OF CONTENTS
# ═══════════════════════════════════════════
def s02_toc(prs):
    sl = blank(prs)
    bg(sl, WHITE)

    # Bottom teal accent bar
    R(sl, 0, SH - Inches(0.08), SW, Inches(0.08), fill=TEAL)
    # Left teal vertical bar (small)
    R(sl, 0, SH - Inches(0.9), Inches(0.06), Inches(0.82), fill=TEAL)

    # "TABLE OF CONTENTS" title (left-bottom area)
    ML(sl,
       l=Inches(0.5), t=Inches(3.6), w=Inches(4.0), h=Inches(2.8),
       items=[
           ('TABLE OF', 36, True, BLACK),
           ('CONTENTS', 36, True, BLACK),
       ],
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    toc_items = [
        ('01', '프로젝트 주제 및 목적'),
        ('02', '서비스 주요 화면'),
        ('03', '기술 스택'),
        ('04', 'DB 구조 설명'),
        ('05', '데이터 등록 기능'),
        ('06', '데이터 조회 기능'),
        ('07', '수정 · 삭제 기능'),
        ('08', '검색 · 필터 기능'),
        ('09', '통계 · 분석 기능'),
        ('10', 'API 테스트 화면'),
        ('11', '프론트엔드 · 백엔드 연동'),
        ('12', '팀원별 역할 및 개선 방향'),
    ]

    col_gap = Inches(0.15)
    badge_r = Inches(0.28)
    row_h   = Inches(0.54)
    start_x = Inches(4.5)
    start_y = Inches(0.35)
    col_w   = Inches(4.2)

    for i, (num, label) in enumerate(toc_items):
        col = i % 2
        row = i // 2
        x = start_x + col * (col_w + col_gap)
        y = start_y + row * row_h

        # Pill background
        bg_fill = RGBColor(0xF0, 0xF0, 0xF0) if i % 2 == 0 else WHITE
        RR(sl, x, y, col_w, Inches(0.46),
           fill=bg_fill, line=LGRAY, lw=Pt(0.5))

        # Number badge (teal circle)
        circ(sl, x + badge_r + Inches(0.06), y + Inches(0.23),
             badge_r, fill=TEAL, text=num, ts=13, tb=True)

        # Label text
        T(sl,
          l=x + badge_r * 2 + Inches(0.16), t=y + Inches(0.06),
          w=col_w - badge_r * 2 - Inches(0.25), h=Inches(0.36),
          text=label, size=13, bold=False, color=DGRAY,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════
# SLIDE 03 — PROJECT OVERVIEW
# ═══════════════════════════════════════════
def s03_overview(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'PROJECT OVERVIEW', '프로젝트 주제 및 목적')

    # Left: decorative circle (placeholder for screenshot)
    circ(sl, cx=Inches(2.8), cy=Inches(4.5), r=Inches(1.9),
         fill=LGRAY, text='', tc=GRAY)
    T(sl, l=Inches(0.9), t=Inches(4.1), w=Inches(3.8), h=Inches(0.8),
      text='앱 화면', size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # Right: description
    content_l = Inches(5.3)
    ML(sl,
       l=content_l, t=Inches(1.85), w=Inches(7.5), h=Inches(1.0),
       items=[
           ('하루온은  목표 수립부터 10분 단위 실행까지,', 16, False, DGRAY),
           ('AI가 함께하는 개인 플래너 웹 서비스입니다.', 16, True, BLACK),
       ],
       align=PP_ALIGN.LEFT, spc=2)

    rows = [
        ('문제', '구글·네이버 캘린더에 없는 연력(年曆) 부재\n목표와 일상 실행의 단절'),
        ('목적', '연→월→주→일→10분 다중 스케일 계획\nAI 인사이트로 사용자 편의 극대화'),
        ('대상', '목표 관리와 일정 실행을 하나의 앱에서\n관리하고 싶은 사용자'),
    ]
    icons = ['⚠', '🎯', '👤']
    for idx, (label, desc) in enumerate(rows):
        y = Inches(2.95) + idx * Inches(1.3)
        circ(sl, cx=content_l + Inches(0.3), cy=y + Inches(0.28),
             r=Inches(0.26), fill=TEAL, text=icons[idx], ts=12, tb=False)
        T(sl, l=content_l + Inches(0.7), t=y,
          w=Inches(6.8), h=Inches(0.32),
          text=label, size=13, bold=True, color=TEAL)
        T(sl, l=content_l + Inches(0.7), t=y + Inches(0.3),
          w=Inches(6.8), h=Inches(0.7),
          text=desc, size=13, color=DGRAY, wrap=True)

    # Keyword pills
    kws = ['#연력', '#10분플래너', '#AI어시스턴트']
    px = content_l
    for kw in kws:
        pill(sl, l=px, t=Inches(6.7), w=Inches(1.7), h=Inches(0.42),
             fill=TEAL, text=kw, ts=12)
        px += Inches(1.85)

# ═══════════════════════════════════════════
# SLIDE 04 — KEY FEATURES
# ═══════════════════════════════════════════
def s04_features(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'KEY FEATURES', '서비스 주요 기능 소개')

    items = [
        ('1', '연력 (Year View)',
         '1년 52주를 한 화면에\n구글·네이버에 없는 차별화 기능'),
        ('2', '캘린더 + 10분 플래너',
         '월간 일정과 집중 실행을\n한 화면에서 관리'),
        ('3', '만다라트 목표 분해',
         '핵심 목표를 9×9 그리드로\n구체적 행동으로 분해'),
        ('4', 'AI 하루온봇',
         '일정 자동 파싱·주간 리캡\n만다라트 자동 생성 제공'),
    ]

    box_w = Inches(2.7)
    box_h = Inches(4.2)
    gap   = Inches(0.38)
    total = 4 * box_w + 3 * gap
    start_x = (SW - total) / 2
    start_y = Inches(1.8)

    for idx, (num, title, desc) in enumerate(items):
        x = start_x + idx * (box_w + gap)

        # Card
        card(sl, l=x, t=start_y, w=box_w, h=box_h, fill=WHITE)

        # Number circle (teal or badge)
        fill_c = TEAL if idx == len(items) - 1 else BADGE
        circ(sl,
             cx=x + box_w / 2, cy=start_y + Inches(0.8),
             r=Inches(0.45),
             fill=fill_c, text=num, ts=20, tb=True)

        # Process label
        pill(sl,
             l=x + box_w / 2 - Inches(0.65), t=start_y + Inches(1.5),
             w=Inches(1.3), h=Inches(0.3),
             fill=TEAL if idx == len(items) - 1 else BADGE,
             text=f'FEATURE {num}', ts=10)

        # Title
        T(sl, l=x + Inches(0.1), t=start_y + Inches(1.95),
          w=box_w - Inches(0.2), h=Inches(0.55),
          text=title, size=15, bold=True, color=BLACK,
          align=PP_ALIGN.CENTER, wrap=True)

        # Description
        T(sl, l=x + Inches(0.15), t=start_y + Inches(2.6),
          w=box_w - Inches(0.3), h=Inches(1.4),
          text=desc, size=12, color=GRAY,
          align=PP_ALIGN.CENTER, wrap=True)

# ═══════════════════════════════════════════
# SLIDE 05 — TECH STACK
# ═══════════════════════════════════════════
def s05_tech(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'TECH STACK', '기술 스택')

    stacks = [
        ('Frontend', 'React 18 + TypeScript\nTailwind CSS + shadcn/ui'),
        ('State Mgmt', 'Zustand + localStorage\n실시간 상태 & 캐싱'),
        ('Database', 'Firebase Firestore\nNoSQL 실시간 동기화'),
        ('Auth', 'Firebase Authentication\nGoogle · Email · Kakao'),
        ('API / AI', 'Vercel Edge Functions\nOpenAI GPT-4o-mini'),
        ('Deploy', 'Vercel\nCI/CD 자동화 배포'),
    ]

    cw = Inches(3.7)
    ch = Inches(2.2)
    gx = Inches(0.35)
    gy = Inches(0.35)
    start_x = Inches(0.55)
    start_y = Inches(1.85)

    for idx, (title, desc) in enumerate(stacks):
        col = idx % 3
        row = idx // 3
        x = start_x + col * (cw + gx)
        y = start_y + row * (ch + gy)

        card(sl, l=x, t=y, w=cw, h=ch, fill=WHITE)

        # Colored top bar on card
        fill_c = TEAL if idx % 2 == 0 else RGBColor(0x2A, 0x2A, 0x2A)
        R(sl, l=x, t=y, w=cw, h=Inches(0.08), fill=fill_c)

        T(sl, l=x + Inches(0.2), t=y + Inches(0.18),
          w=cw - Inches(0.4), h=Inches(0.42),
          text=title, size=16, bold=True, color=BLACK)

        T(sl, l=x + Inches(0.2), t=y + Inches(0.65),
          w=cw - Inches(0.4), h=Inches(1.3),
          text=desc, size=13, color=DGRAY, wrap=True)

# ═══════════════════════════════════════════
# SLIDE 06 — DATABASE STRUCTURE
# ═══════════════════════════════════════════
def s06_db(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'DATABASE STRUCTURE', 'DB 컬렉션 구조')

    # Left: collection tree
    lx = Inches(0.5)
    card(sl, l=lx, t=Inches(1.75), w=Inches(5.8), h=Inches(5.4), fill=WHITE)

    T(sl, l=lx + Inches(0.2), t=Inches(1.9),
      w=Inches(5.4), h=Inches(0.4),
      text='Firestore 컬렉션 구조', size=14, bold=True, color=TEAL)

    tree_lines = [
        'users/{uid}/',
        '  ├─ profile     uid · email · displayName',
        '  ├─ events      id · date · title · color',
        '  ├─ todos/{date}  날짜별 할 일',
        '  ├─ mandala     cells[81] · updatedAt',
        '  └─ diaries/{date}  날짜별 일기',
        '',
        'shares/{shareId}',
        '  ownerUid · scope · payload',
        '  expiresAt (30일 TTL)',
    ]
    ML(sl,
       l=lx + Inches(0.2), t=Inches(2.35),
       w=Inches(5.4), h=Inches(4.5),
       items=tree_lines,
       default_size=12, default_color=DGRAY,
       wrap=False, spc=2)

    # Right: events fields table
    rx = Inches(6.8)
    card(sl, l=rx, t=Inches(1.75), w=Inches(6.0), h=Inches(5.4), fill=WHITE)

    T(sl, l=rx + Inches(0.2), t=Inches(1.9),
      w=Inches(5.6), h=Inches(0.4),
      text='events 컬렉션 주요 필드', size=14, bold=True, color=TEAL)

    fields = [
        ('필드명',     '타입',    '설명',         True),
        ('id',        'string',  '이벤트 고유 ID',False),
        ('date',      'string',  'YYYY-MM-DD',   False),
        ('title',     'string',  '일정 제목',     False),
        ('color',     'string',  '색상 코드(#hex)',False),
        ('startTime', 'string',  'HH:MM (선택)',  False),
        ('endTime',   'string',  'HH:MM (선택)',  False),
        ('createdBy', 'string',  '작성자 uid',    False),
        ('createdAt', 'number',  '생성 타임스탬프',False),
    ]
    row_h = Inches(0.44)
    ty = Inches(2.38)
    for r_idx, (f, tp, desc, is_hdr) in enumerate(fields):
        row_y = ty + r_idx * row_h
        row_fill = RGBColor(0xE8, 0xFA, 0xF6) if is_hdr else (
            RGBColor(0xF8, 0xF8, 0xF8) if r_idx % 2 == 0 else WHITE
        )
        R(sl, l=rx + Inches(0.1), t=row_y,
          w=Inches(5.8), h=row_h,
          fill=TEAL if is_hdr else row_fill)

        col_cfg = [
            (Inches(0.1),  Inches(1.3)),
            (Inches(1.5),  Inches(1.1)),
            (Inches(2.7),  Inches(2.9)),
        ]
        for val, (off_x, col_w) in zip([f, tp, desc], col_cfg):
            T(sl,
              l=rx + Inches(0.1) + off_x,
              t=row_y + Inches(0.05),
              w=col_w, h=row_h - Inches(0.08),
              text=val,
              size=11 if not is_hdr else 12,
              bold=is_hdr,
              color=WHITE if is_hdr else DGRAY,
              anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════
# SLIDE 07 — CREATE (데이터 등록)
# ═══════════════════════════════════════════
def s07_create(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'CREATE', '데이터 등록 기능 시연')

    # Left placeholder
    card(sl, l=Inches(0.4), t=Inches(1.75), w=Inches(5.5), h=Inches(5.35), fill=WHITE)
    T(sl, l=Inches(0.4), t=Inches(1.75), w=Inches(5.5), h=Inches(5.35),
      text='📱  화면 캡처\n(이벤트 등록 / 할 일 추가 / 일기 작성)',
      size=15, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True)

    # Right: step cards
    steps = [
        ('STEP 1', '이벤트 등록',
         '캘린더에서 날짜 클릭 →\n폼 입력 → Firestore 저장',
         'addEvent()'),
        ('STEP 2', '할 일 추가',
         '오늘 뷰에서 텍스트 입력 →\n날짜별 Todo 저장',
         'upsertTodosForDate()'),
        ('STEP 3', '일기 작성',
         '일기 탭에서 내용 입력 →\n날짜별 Diary 저장',
         'upsertDiary()'),
    ]
    rx = Inches(6.3)
    for i, (badge, title, desc, fn) in enumerate(steps):
        y = Inches(1.85) + i * Inches(1.75)
        card(sl, l=rx, t=y, w=Inches(6.6), h=Inches(1.58), fill=WHITE)
        pill(sl, l=rx + Inches(0.15), t=y + Inches(0.15),
             w=Inches(0.95), h=Inches(0.3),
             fill=TEAL, text=badge, ts=10)
        T(sl, l=rx + Inches(1.2), t=y + Inches(0.1),
          w=Inches(5.2), h=Inches(0.38),
          text=title, size=15, bold=True, color=BLACK)
        T(sl, l=rx + Inches(0.18), t=y + Inches(0.5),
          w=Inches(4.5), h=Inches(0.7),
          text=desc, size=12, color=DGRAY, wrap=True)
        T(sl, l=rx + Inches(4.7), t=y + Inches(0.5),
          w=Inches(1.7), h=Inches(0.55),
          text=fn, size=10, color=TEAL, italic=True, wrap=False)

# ═══════════════════════════════════════════
# SLIDE 08 — READ (데이터 조회)
# ═══════════════════════════════════════════
def s08_read(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'READ', '데이터 조회 기능 시연')

    views = [
        ('연력 뷰', '1년 52주 전체 활동을\n한눈에 조회', 'fetchEvents()'),
        ('월간 캘린더', '이벤트 컬러바로\n일정 시각화', 'fetchEvents()'),
        ('오늘 뷰', '당일 일정 + 할 일\n목록 통합 조회', 'fetchTodos()'),
    ]

    cw = Inches(3.7)
    ch = Inches(5.0)
    gap = Inches(0.38)
    total = 3 * cw + 2 * gap
    sx = (SW - total) / 2

    for i, (title, desc, fn) in enumerate(views):
        x = sx + i * (cw + gap)
        card(sl, l=x, t=Inches(1.8), w=cw, h=ch, fill=WHITE)

        # Placeholder for screenshot
        R(sl, l=x + Inches(0.15), t=Inches(2.0),
          w=cw - Inches(0.3), h=Inches(2.8),
          fill=RGBColor(0xEA, 0xEA, 0xEA))
        T(sl,
          l=x + Inches(0.15), t=Inches(2.0),
          w=cw - Inches(0.3), h=Inches(2.8),
          text='화면 캡처', size=14, color=GRAY,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        T(sl, l=x + Inches(0.15), t=Inches(4.9),
          w=cw - Inches(0.3), h=Inches(0.42),
          text=title, size=15, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        T(sl, l=x + Inches(0.15), t=Inches(5.35),
          w=cw - Inches(0.3), h=Inches(0.7),
          text=desc, size=12, color=GRAY,
          align=PP_ALIGN.CENTER, wrap=True)
        T(sl, l=x + Inches(0.15), t=Inches(6.1),
          w=cw - Inches(0.3), h=Inches(0.35),
          text=fn, size=11, color=TEAL, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 09 — UPDATE & DELETE
# ═══════════════════════════════════════════
def s09_update(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'UPDATE & DELETE', '수정 · 삭제 기능 시연')

    # Left card (Update)
    lx = Inches(0.5)
    card(sl, l=lx, t=Inches(1.75), w=Inches(5.8), h=Inches(5.35), fill=WHITE)
    pill(sl, l=lx + Inches(0.2), t=Inches(1.95),
         w=Inches(1.2), h=Inches(0.32), fill=TEAL, text='UPDATE')
    T(sl, l=lx + Inches(0.2), t=Inches(2.42),
      w=Inches(5.4), h=Inches(0.4),
      text='이벤트 수정', size=18, bold=True, color=BLACK)
    update_steps = [
        '① 이벤트 클릭 → 수정 팝업 열기',
        '② 제목 · 시간 · 색상 변경',
        '③ 저장 버튼 → Firestore 업데이트',
        '',
        '함수: updateEvent(id, patch)',
    ]
    ML(sl, l=lx + Inches(0.2), t=Inches(3.0),
       w=Inches(5.4), h=Inches(2.5),
       items=update_steps,
       default_size=13, default_color=DGRAY, spc=3)

    # Right card (Delete)
    rx = Inches(6.9)
    card(sl, l=rx, t=Inches(1.75), w=Inches(5.8), h=Inches(5.35), fill=WHITE)
    pill(sl, l=rx + Inches(0.2), t=Inches(1.95),
         w=Inches(1.2), h=Inches(0.32),
         fill=BADGE, text='DELETE', tc=WHITE)
    T(sl, l=rx + Inches(0.2), t=Inches(2.42),
      w=Inches(5.4), h=Inches(0.4),
      text='이벤트 · 할 일 삭제', size=18, bold=True, color=BLACK)
    delete_steps = [
        '① 이벤트 삭제 버튼 클릭',
        '② 확인 다이얼로그 표시',
        '③ 확인 → Firestore 문서 삭제',
        '',
        '함수: removeEvent(id)',
        '함수: deleteHighlight(uid, id)',
    ]
    ML(sl, l=rx + Inches(0.2), t=Inches(3.0),
       w=Inches(5.4), h=Inches(2.8),
       items=delete_steps,
       default_size=13, default_color=DGRAY, spc=3)

    # Center illustration area
    T(sl, l=Inches(0.5), t=Inches(5.5), w=Inches(12.0), h=Inches(0.5),
      text='📸  수정 팝업 / 삭제 확인 화면 캡처 삽입',
      size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════
# SLIDE 10 — SEARCH & FILTER
# ═══════════════════════════════════════════
def s10_search(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'SEARCH & FILTER', '검색 · 필터 기능 시연')

    # Center screenshot area
    card(sl, l=Inches(0.4), t=Inches(1.75), w=Inches(7.5), h=Inches(5.35), fill=WHITE)
    T(sl, l=Inches(0.4), t=Inches(1.75), w=Inches(7.5), h=Inches(5.35),
      text='📱  검색 / 필터 화면 캡처\n(날짜별 조회 · 기간별 탐색)',
      size=15, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True)

    # Right: feature list
    rx = Inches(8.3)
    features = [
        ('날짜별 필터', '캘린더 날짜 클릭 →\n해당 날짜 이벤트만 표시'),
        ('기간별 조회', '주간 뷰에서 날짜 범위\n슬라이드로 탐색'),
        ('만다라트 조회', '목표별 하위 항목\n그리드에서 확인'),
        ('API 엔드포인트', 'GET /api/events\n?uid=...&date=YYYY-MM-DD'),
    ]
    for i, (title, desc) in enumerate(features):
        y = Inches(1.85) + i * Inches(1.25)
        card(sl, l=rx, t=y, w=Inches(4.6), h=Inches(1.1), fill=WHITE)
        circ(sl, cx=rx + Inches(0.35), cy=y + Inches(0.55),
             r=Inches(0.26), fill=TEAL, text=str(i + 1), ts=13)
        T(sl, l=rx + Inches(0.75), t=y + Inches(0.1),
          w=Inches(3.7), h=Inches(0.38),
          text=title, size=14, bold=True, color=BLACK)
        T(sl, l=rx + Inches(0.75), t=y + Inches(0.48),
          w=Inches(3.7), h=Inches(0.55),
          text=desc, size=11, color=GRAY, wrap=True)

# ═══════════════════════════════════════════
# SLIDE 11 — STATISTICS
# ═══════════════════════════════════════════
def s11_stats(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'STATISTICS', '통계 · 분석 기능 시연')

    stats = [
        ('TRACTION 1', '할 일 완료율',
         '완료된 할 일 / 전체 할 일\n오늘 뷰에서 실시간 집계'),
        ('TRACTION 2', 'AI 주간 리캡',
         '한 주 요약 + 하이라이트\n+ AI 제안 자동 생성'),
        ('TRACTION 3', '연간 달성 현황',
         '연력에서 색상 농도로\n활동 밀도 시각화'),
    ]

    cw = Inches(3.7)
    ch = Inches(4.5)
    gap = Inches(0.38)
    total = 3 * cw + 2 * gap
    sx = (SW - total) / 2

    for i, (badge, title, desc) in enumerate(stats):
        x = sx + i * (cw + gap)
        card(sl, l=x, t=Inches(1.85), w=cw, h=ch, fill=WHITE)

        # Badge
        pill(sl, l=x + cw / 2 - Inches(0.8), t=Inches(2.0),
             w=Inches(1.6), h=Inches(0.3),
             fill=TEAL, text=badge, ts=10)

        # Chart placeholder
        R(sl, l=x + Inches(0.2), t=Inches(2.45),
          w=cw - Inches(0.4), h=Inches(2.1),
          fill=RGBColor(0xEA, 0xEA, 0xEA))
        T(sl, l=x + Inches(0.2), t=Inches(2.45),
          w=cw - Inches(0.4), h=Inches(2.1),
          text='차트 / 화면', size=13, color=GRAY,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Title
        T(sl, l=x + Inches(0.15), t=Inches(4.65),
          w=cw - Inches(0.3), h=Inches(0.45),
          text=title, size=15, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

        # Description
        T(sl, l=x + Inches(0.15), t=Inches(5.15),
          w=cw - Inches(0.3), h=Inches(0.9),
          text=desc, size=12, color=GRAY,
          align=PP_ALIGN.CENTER, wrap=True)

# ═══════════════════════════════════════════
# SLIDE 12 — API TEST
# ═══════════════════════════════════════════
def s12_api(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'API TEST', 'API 테스트 화면 (Swagger 대체)')

    # Left: endpoint cards
    endpoints = [
        ('POST', '/api/ai/insight',     '일일 AI 인사이트 생성'),
        ('POST', '/api/ai/parse-event', '자연어 → 이벤트 파싱'),
        ('POST', '/api/ai/mandala',     '목표 → 만다라트 자동 분해'),
        ('POST', '/api/ai/recap',       '주간 요약 자동 생성'),
    ]
    lx = Inches(0.4)
    for i, (method, ep, desc) in enumerate(endpoints):
        y = Inches(1.8) + i * Inches(1.32)
        card(sl, l=lx, t=y, w=Inches(6.0), h=Inches(1.15), fill=WHITE)
        pill(sl, l=lx + Inches(0.15), t=y + Inches(0.4),
             w=Inches(0.85), h=Inches(0.32),
             fill=TEAL, text=method, ts=11)
        T(sl, l=lx + Inches(1.1), t=y + Inches(0.12),
          w=Inches(4.7), h=Inches(0.38),
          text=ep, size=14, bold=True, color=BLACK)
        T(sl, l=lx + Inches(1.1), t=y + Inches(0.55),
          w=Inches(4.7), h=Inches(0.38),
          text=desc, size=12, color=GRAY)

    # Right: description
    rx = Inches(7.0)
    card(sl, l=rx, t=Inches(1.8), w=Inches(5.8), h=Inches(5.3), fill=WHITE)
    T(sl, l=rx + Inches(0.2), t=Inches(2.0),
      w=Inches(5.4), h=Inches(0.45),
      text='API 동작 확인 방법', size=16, bold=True, color=TEAL)

    notes = [
        '✅ Vercel 대시보드에서 함수 호출 로그 확인',
        '✅ Firebase Firestore 콘솔에서\n   데이터 저장/변경 실시간 확인',
        '✅ 브라우저 개발자 도구 Network 탭에서\n   요청/응답 상세 내용 확인',
        '',
        '⚙  FastAPI Swagger(/docs) 대신\n   Vercel Edge Functions 방식으로\n   API 서버 구현',
    ]
    ML(sl, l=rx + Inches(0.2), t=Inches(2.6),
       w=Inches(5.4), h=Inches(4.2),
       items=notes,
       default_size=13, default_color=DGRAY,
       wrap=True, spc=4)

# ═══════════════════════════════════════════
# SLIDE 13 — ARCHITECTURE
# ═══════════════════════════════════════════
def s13_arch(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'ARCHITECTURE', '프론트엔드 · 백엔드 연동 구조')

    steps = [
        ('1', '사용자 입력',
         'React 화면에서\n이벤트 · 할 일 · 일기 입력'),
        ('2', '상태 관리',
         'Zustand Store ↔ localStorage\n캐싱으로 오프라인 대응'),
        ('3', 'Firebase 저장',
         'Firebase SDK로 Firestore 동기화\n보안 규칙: 소유자만 접근'),
        ('4', 'AI 처리',
         'Vercel Edge Function 호출\nToken 검증 → GPT-4o-mini'),
    ]

    cw = Inches(2.7)
    ch = Inches(4.3)
    gap = Inches(0.35)
    total = 4 * cw + 3 * gap
    sx = (SW - total) / 2
    sy = Inches(1.9)

    for i, (num, title, desc) in enumerate(steps):
        x = sx + i * (cw + gap)
        card(sl, l=x, t=sy, w=cw, h=ch, fill=WHITE)

        # Teal header bar
        R(sl, l=x, t=sy, w=cw, h=Inches(0.08), fill=TEAL)

        # Number circle
        circ(sl, cx=x + cw / 2, cy=sy + Inches(0.75),
             r=Inches(0.42), fill=TEAL, text=num, ts=20)

        # Process label
        pill(sl, l=x + cw / 2 - Inches(0.65), t=sy + Inches(1.35),
             w=Inches(1.3), h=Inches(0.3),
             fill=BADGE, text=f'PROCESS {num}', ts=10)

        T(sl, l=x + Inches(0.1), t=sy + Inches(1.8),
          w=cw - Inches(0.2), h=Inches(0.45),
          text=title, size=15, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        T(sl, l=x + Inches(0.1), t=sy + Inches(2.35),
          w=cw - Inches(0.2), h=Inches(1.7),
          text=desc, size=12, color=GRAY,
          align=PP_ALIGN.CENTER, wrap=True)

        # Arrow between cards
        if i < len(steps) - 1:
            ax = x + cw + gap / 2 - Inches(0.1)
            T(sl, l=ax, t=sy + ch / 2 - Inches(0.2),
              w=Inches(0.3), h=Inches(0.4),
              text='→', size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Bottom summary bar
    R(sl, l=Inches(0.4), t=Inches(6.55),
      w=SW - Inches(0.8), h=Inches(0.55),
      fill=TEAL)
    T(sl, l=Inches(0.4), t=Inches(6.55),
      w=SW - Inches(0.8), h=Inches(0.55),
      text='React → Zustand → Firebase Firestore  +  React → Vercel Edge Functions → OpenAI API',
      size=13, bold=True, color=WHITE,
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════
# SLIDE 14 — OUR TEAM
# ═══════════════════════════════════════════
def s14_team(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'OUR TEAM', '팀원별 역할')

    members = [
        ('김태동', 'APP DEV', '모바일 UI\nReact 컴포넌트\nPWA 설정'),
        ('박성훈', 'BACKEND', 'Firebase 구조 설계\nFirestore 스키마\n보안 규칙 작성'),
        ('이현아', 'PLANNER\n& WEB DEV', '서비스 기획\n웹 반응형 구현\n연력 · 플래너 뷰'),
        ('이창민', 'AI / BOT', 'Vercel Edge Functions\nOpenAI API 연동\n하루온봇 개발'),
    ]

    mw = Inches(2.75)
    gap = Inches(0.38)
    total = 4 * mw + 3 * gap
    sx = (SW - total) / 2
    sy = Inches(1.75)

    # Connector line
    line_y = sy + Inches(1.05)
    R(sl, l=sx + Inches(0.3), t=line_y,
      w=total - Inches(0.6), h=Inches(0.025),
      fill=LGRAY)

    for i, (name, role, detail) in enumerate(members):
        x = sx + i * (mw + gap)

        # Profile circle (teal for highlight, gray for others)
        is_highlight = (i == 2)  # 이현아 highlighted
        fill_c = TEAL if is_highlight else LGRAY
        circ(sl, cx=x + mw / 2, cy=sy + Inches(0.9),
             r=Inches(0.9), fill=fill_c, text='', tc=WHITE)
        T(sl, l=x, t=sy + Inches(0.15),
          w=mw, h=Inches(1.5),
          text='프로필', size=13, color=WHITE if is_highlight else GRAY,
          align=PP_ALIGN.CENTER)

        # Name
        T(sl, l=x, t=sy + Inches(2.1),
          w=mw, h=Inches(0.48),
          text=name, size=18, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

        # Role badge
        pill(sl, l=x + mw / 2 - Inches(0.8), t=sy + Inches(2.65),
             w=Inches(1.6), h=Inches(0.3),
             fill=TEAL if is_highlight else BADGE,
             text=role.replace('\n', ' '), ts=10)

        # Detail
        T(sl, l=x + Inches(0.1), t=sy + Inches(3.1),
          w=mw - Inches(0.2), h=Inches(1.8),
          text=detail, size=12, color=DGRAY,
          align=PP_ALIGN.CENTER, wrap=True)

    T(sl, l=0, t=Inches(7.0), w=SW, h=Inches(0.35),
      text='* 프로필 사진은 별도 삽입 예정',
      size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════
# SLIDE 15 — ROADMAP
# ═══════════════════════════════════════════
def s15_roadmap(prs):
    sl = blank(prs)
    bg(sl, BG)
    hdr(sl, 'ROADMAP', '개선 방향')

    milestones = [
        ('단기', 'AI 기능 완성',
         '목표\nAI 4가지 기능 구현 완료\n\n전략\ninsight · parse-event\nmandala · recap API\n개발 및 연동'),
        ('중기', '협업 기능 강화',
         '목표\n팀 플랜 공유 고도화\n\n전략\n실시간 협업 모드\n공유 링크 TTL 연장\n팀원 초대 기능'),
        ('장기', 'PWA 앱 출시',
         '목표\n모바일 앱 & 오프라인 지원\n\n전략\nService Worker 캐싱\n푸시 알림\n소셜 플래너 기능'),
    ]

    cw = Inches(3.7)
    ch = Inches(5.0)
    gap = Inches(0.38)
    total = 3 * cw + 2 * gap
    sx = (SW - total) / 2

    # Timeline connector
    R(sl, l=sx + cw / 2, t=Inches(2.55),
      w=total - cw, h=Inches(0.04), fill=LGRAY)

    for i, (period, title, detail) in enumerate(milestones):
        x = sx + i * (cw + gap)
        is_highlight = (i == 0)

        # Milestone circle on timeline
        circ(sl, cx=x + cw / 2, cy=Inches(2.57),
             r=Inches(0.4),
             fill=TEAL if is_highlight else BADGE,
             text='', tc=WHITE)

        # Card below circle
        card(sl, l=x, t=Inches(3.05),
             w=cw, h=ch - Inches(0.85),
             fill=TEAL if is_highlight else WHITE,
             line=None if is_highlight else LGRAY)

        # Period & title (above card)
        T(sl, l=x, t=Inches(1.85), w=cw, h=Inches(0.35),
          text=period, size=14, bold=True, color=TEAL,
          align=PP_ALIGN.CENTER)
        T(sl, l=x, t=Inches(2.2), w=cw, h=Inches(0.35),
          text=title, size=14, bold=True, color=BLACK,
          align=PP_ALIGN.CENTER)

        # Detail lines inside card
        for j, line in enumerate(detail.split('\n')):
            is_label = line in ('목표', '전략')
            T(sl, l=x + Inches(0.2), t=Inches(3.15) + j * Inches(0.37),
              w=cw - Inches(0.4), h=Inches(0.35),
              text=line,
              size=12 if not is_label else 13,
              bold=is_label,
              color=(WHITE if is_highlight else
                     (TEAL if is_label else DGRAY)))

# ═══════════════════════════════════════════
# SLIDE 16 — CLOSING
# ═══════════════════════════════════════════
def s16_closing(prs):
    sl = blank(prs)
    bg(sl, RGBColor(0xEB, 0xEB, 0xEB))

    # Same arc decorations as cover
    OV(sl, l=Inches(8.5), t=Inches(-1.0),
       w=Inches(6.5), h=Inches(6.5),
       fill=None, line=TEAL, lw=Pt(8))
    OV(sl, l=Inches(-1.5), t=Inches(2.5),
       w=Inches(6), h=Inches(6),
       fill=None, line=LGRAY, lw=Pt(8))

    # Center card
    card(sl, l=Inches(2.8), t=Inches(1.1),
         w=Inches(7.7), h=Inches(5.3),
         fill=WHITE, line=LGRAY)

    T(sl, l=Inches(2.8), t=Inches(1.9), w=Inches(7.7), h=Inches(0.5),
      text='THANK YOU', size=13, bold=False, color=GRAY,
      align=PP_ALIGN.CENTER)

    T(sl, l=Inches(2.8), t=Inches(2.4), w=Inches(7.7), h=Inches(1.0),
      text='"연간 계획부터 10분 실행까지,\nAI와 함께 채우는 하루온"',
      size=20, bold=True, color=BLACK,
      align=PP_ALIGN.CENTER, wrap=True)

    T(sl, l=Inches(2.8), t=Inches(3.7), w=Inches(7.7), h=Inches(0.4),
      text='GitHub: github.com/halee0426/KDT_webai_team_planner',
      size=13, color=TEAL, align=PP_ALIGN.CENTER)

    T(sl, l=Inches(2.8), t=Inches(4.4), w=Inches(7.7), h=Inches(0.45),
      text='KDT WebAI 4팀  |  김태동 · 박성훈 · 이현아 · 이창민',
      size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════
def main():
    prs = new_prs()

    print("Generating slides...")
    s01_cover(prs);    print("  [OK] Slide 01 - Cover")
    s02_toc(prs);      print("  [OK] Slide 02 - TOC")
    s03_overview(prs); print("  [OK] Slide 03 - Project Overview")
    s04_features(prs); print("  [OK] Slide 04 - Key Features")
    s05_tech(prs);     print("  [OK] Slide 05 - Tech Stack")
    s06_db(prs);       print("  [OK] Slide 06 - DB Structure")
    s07_create(prs);   print("  [OK] Slide 07 - Create")
    s08_read(prs);     print("  [OK] Slide 08 - Read")
    s09_update(prs);   print("  [OK] Slide 09 - Update & Delete")
    s10_search(prs);   print("  [OK] Slide 10 - Search & Filter")
    s11_stats(prs);    print("  [OK] Slide 11 - Statistics")
    s12_api(prs);      print("  [OK] Slide 12 - API Test")
    s13_arch(prs);     print("  [OK] Slide 13 - Architecture")
    s14_team(prs);     print("  [OK] Slide 14 - Team")
    s15_roadmap(prs);  print("  [OK] Slide 15 - Roadmap")
    s16_closing(prs);  print("  [OK] Slide 16 - Closing")

    out = "HaruOn_PPT.pptx"
    prs.save(out)
    print("\nSaved: " + out)

if __name__ == "__main__":
    main()
